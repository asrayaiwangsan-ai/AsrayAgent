
from io import BytesIO
import numpy as np
from docx import Document, ImagePart
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table, _Cell
from docx.text.paragraph import Paragraph
from PIL import Image
from rapidocr_onnxruntime import RapidOCR
import uuid
import fitz  # pyMuPDF里面的fitz包，不要与pip install fitz混淆
import numpy as np

def doc2text(filepath, multimodal, base_dir):

    ocr = RapidOCR()
    doc = Document(filepath)
    resp = ""
    pics = []

    ret = []

    def iter_block_items(parent):
        from docx.document import Document
        if isinstance(parent, Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("RapidOCRDocLoader parse fail")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    for i, block in enumerate(iter_block_items(doc)):
        if isinstance(block, Paragraph):
            resp += block.text.strip() + "\n"

            images = block._element.xpath(".//pic:pic")  # 获取所有图片
            for image in images:
                for img_id in image.xpath(".//a:blip/@r:embed"):  # 获取图片id
                    part = doc.part.related_parts[
                        img_id
                    ]
                    if isinstance(part, ImagePart):
                        image = Image.open(BytesIO(part._blob))
                        if (multimodal):
                            name = base_dir + "/images/" + uuid.uuid4().__str__() + ".png"
                            resp = resp + "<image>"
                            image.save(name)
                            pics.append(name)
                        else:
                            image = Image.open(BytesIO(part._blob))
                            result, _ = ocr(np.array(image))
                            if result:
                                ocr_result = [line[1] for line in result]
                                ret.append(
                                    {"type": "text", "text":  "\n".join(ocr_result)})

    return resp, pics

def pdf2text(filepath, multimodal, base_dir):
    
    ocr = RapidOCR()
    doc = fitz.open(filepath)
    resp = ""
    pics = []

    for i, page in enumerate(doc):
        text = page.get_text("text")
        resp += text + "\n"

        img_list = page.get_image_info(xrefs=True)
        for img in img_list:
            if xref := img.get("xref"):
                pix = fitz.Pixmap(doc, xref)
                img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)
                if img_array.shape[2] == 1:
                    img_array = img_array[:, :, 0]

                if (multimodal):
                    image = Image.fromarray(img_array)
                    name = base_dir + "/images/" + uuid.uuid4().__str__() + ".png"
                    resp = resp + "<image>"
                    image.save(name)
                    pics.append(name)

                else:
                    result, _ = ocr(img_array)
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                    
                    
    return resp, pics


def pptx2text(filepath, multimodal, base_dir):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(filepath)
    resp = ""
    pics = []
    ocr = RapidOCR()

    for i, slide in enumerate(prs.slides):
        resp += f"--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            # 提取文本
            if hasattr(shape, "text"):
                resp += shape.text.strip() + "\n"
            
            # 提取图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_bytes = shape.image.blob
                image = Image.open(BytesIO(image_bytes))
                
                if multimodal:
                    name = base_dir + "/images/" + uuid.uuid4().__str__() + ".png"
                    resp += "<image>"
                    image.save(name)
                    pics.append(name)
                else:
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result) + "\n"
    
    return resp, pics


def xlsx2text(filepath, multimodal, base_dir):
    import pandas as pd
    
    resp = ""
    try:
        # 读取所有 sheet
        excel_data = pd.read_excel(filepath, sheet_name=None)
        for sheet_name, df in excel_data.items():
            resp += f"--- Sheet: {sheet_name} ---\n"
            # 转换为 Markdown 表格格式，这对 LLM 非常友好
            resp += df.to_markdown(index=False) + "\n\n"
    except Exception as e:
        resp = f"解析 Excel 失败: {str(e)}"
    
    return resp, []
